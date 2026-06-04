from dataclasses import dataclass
from pathlib import Path
from tempfile import NamedTemporaryFile

import fitz
from docx import Document as DocxDocument
from pptx import Presentation

from app.core.config import settings
from app.models.document import Document
from app.services.storage_service import StorageService


@dataclass(frozen=True)
class ParsedBlock:
    text: str
    page_number: int | None
    source_type: str


@dataclass(frozen=True)
class TextChunk:
    content: str
    page_number: int | None
    chunk_index: int
    source_type: str
    token_count: int


class IngestionService:
    def __init__(self, storage: StorageService | None = None) -> None:
        self.storage = storage or StorageService()

    def load_blocks(self, document: Document) -> list[ParsedBlock]:
        suffix = Path(document.original_filename).suffix.lower()
        with NamedTemporaryFile(suffix=suffix) as temp_file:
            self.storage.download_file(document.storage_key, temp_file.name)
            if suffix == ".pdf":
                return self._parse_pdf(temp_file.name)
            if suffix == ".pptx":
                return self._parse_pptx(temp_file.name)
            if suffix == ".docx":
                return self._parse_docx(temp_file.name)
        raise ValueError(f"Unsupported document extension: {suffix}")

    def chunk_blocks(self, blocks: list[ParsedBlock]) -> list[TextChunk]:
        chunks: list[TextChunk] = []
        current_parts: list[str] = []
        current_tokens = 0
        current_page: int | None = None
        current_source = "document"

        for block in blocks:
            sentences = self._split_sentences(block.text)
            for sentence in sentences:
                sentence_tokens = self._estimate_tokens(sentence)
                if current_parts and current_tokens + sentence_tokens > settings.chunk_target_tokens:
                    chunks.append(self._build_chunk(chunks, current_parts, current_page, current_source))
                    overlap_text = self._overlap_text(current_parts)
                    current_parts = [overlap_text] if overlap_text else []
                    current_tokens = self._estimate_tokens(overlap_text) if overlap_text else 0

                current_parts.append(sentence)
                current_tokens += sentence_tokens
                current_page = current_page or block.page_number
                current_source = block.source_type

        if current_parts:
            chunks.append(self._build_chunk(chunks, current_parts, current_page, current_source))

        return chunks

    def _parse_pdf(self, path: str) -> list[ParsedBlock]:
        blocks: list[ParsedBlock] = []
        with fitz.open(path) as pdf:
            for page_index, page in enumerate(pdf, start=1):
                text = page.get_text("text").strip()
                if text:
                    blocks.append(ParsedBlock(text=text, page_number=page_index, source_type="pdf"))
        return blocks

    def _parse_pptx(self, path: str) -> list[ParsedBlock]:
        presentation = Presentation(path)
        blocks: list[ParsedBlock] = []
        for slide_index, slide in enumerate(presentation.slides, start=1):
            text_parts: list[str] = []
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    text_parts.append(shape.text.strip())
            if text_parts:
                blocks.append(ParsedBlock(text="\n".join(text_parts), page_number=slide_index, source_type="pptx"))
        return blocks

    def _parse_docx(self, path: str) -> list[ParsedBlock]:
        document = DocxDocument(path)
        blocks = [
            ParsedBlock(text=paragraph.text.strip(), page_number=None, source_type="docx")
            for paragraph in document.paragraphs
            if paragraph.text.strip()
        ]
        return blocks

    def _split_sentences(self, text: str) -> list[str]:
        paragraphs = [part.strip() for part in text.replace("\r\n", "\n").split("\n") if part.strip()]
        return paragraphs or [text.strip()]

    def _estimate_tokens(self, text: str) -> int:
        return max(1, len(text.split()) * 4 // 3)

    def _overlap_text(self, parts: list[str]) -> str:
        words = " ".join(parts).split()
        if len(words) <= settings.chunk_overlap_tokens:
            return ""
        return " ".join(words[-settings.chunk_overlap_tokens :])

    def _build_chunk(
        self,
        existing_chunks: list[TextChunk],
        parts: list[str],
        page_number: int | None,
        source_type: str,
    ) -> TextChunk:
        content = "\n\n".join(part for part in parts if part.strip()).strip()
        return TextChunk(
            content=content,
            page_number=page_number,
            chunk_index=len(existing_chunks),
            source_type=source_type,
            token_count=self._estimate_tokens(content),
        )
